"""
office_extractor.py — Office 文档提取器 (v0.8)

支持格式：Word (.docx), Excel (.xlsx), PowerPoint (.pptx)

输出统一为 DocumentElement 列表，与 pdf_extractor 格式一致。

用法：
    from src.data.extractors.office_extractor import OfficeExtractor
    ext = OfficeExtractor()
    elements = ext.extract("report.docx")
"""

import logging
from pathlib import Path
from typing import List

from src.data.extractors.pdf_extractor import DocumentElement

logger = logging.getLogger(__name__)


class OfficeExtractor:
    """Office 文档提取器。

    支持的格式通过文件扩展名自动识别。
    """

    def extract(self, file_path: str) -> List[DocumentElement]:
        path = Path(file_path)
        ext = path.suffix.lower()
        name = path.name

        if ext in (".docx", ".doc"):
            return self._extract_docx(file_path)
        elif ext in (".xlsx", ".xls"):
            return self._extract_xlsx(file_path)
        elif ext in (".pptx", ".ppt"):
            return self._extract_pptx(file_path)
        else:
            logger.warning("不支持的 Office 格式: %s", ext)
            return []

    # ------------------------------------------------------------------
    # Word
    # ------------------------------------------------------------------

    def _extract_docx(self, file_path: str) -> List[DocumentElement]:
        """提取 Word 文档：按段落输出。"""
        try:
            from docx import Document
        except ImportError:
            logger.warning("python-docx 未安装。pip install python-docx")
            return self._fallback_text(file_path)

        doc = Document(file_path)
        elements = []
        page = 1  # Word 无原生分页，按段落数估算

        for i, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if not text:
                continue

            # 判断元素类型
            if para.style and para.style.name and "Heading" in para.style.name:
                cat = "Title"
            elif any(text.startswith(p) for p in ("- ", "• ", "1.", "a)", "①")):
                cat = "ListItem"
            else:
                cat = "NarrativeText"

            elements.append(DocumentElement(
                category=cat,
                text=text,
                page_number=page + i // 40,  # 约 40 段 = 1 页
                metadata={"source_file": file_path, "paragraph_index": i},
            ))

        # 提取表格
        for i, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                elements.append(DocumentElement(
                    category="Table",
                    text="\n".join(rows),
                    page_number=page,
                    metadata={"source_file": file_path, "table_index": i},
                ))

        logger.info("Word 提取完成: %d 元素 (%s)", len(elements), Path(file_path).name)
        return elements

    # ------------------------------------------------------------------
    # Excel
    # ------------------------------------------------------------------

    def _extract_xlsx(self, file_path: str) -> List[DocumentElement]:
        """提取 Excel：每个 Sheet 作为一个 NarrativeText + 表格。"""
        try:
            import openpyxl
        except ImportError:
            logger.warning("openpyxl 未安装。pip install openpyxl")
            return self._fallback_text(file_path)

        wb = openpyxl.load_workbook(file_path, data_only=True)
        elements = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    row_text = " | ".join(str(c) if c is not None else "" for c in row)
                    rows.append(row_text)

            if rows:
                elements.append(DocumentElement(
                    category="NarrativeText",
                    text=f"Sheet: {sheet_name}",
                    page_number=1,
                    metadata={"source_file": file_path, "sheet": sheet_name},
                ))
                elements.append(DocumentElement(
                    category="Table",
                    text="\n".join(rows[:200]),  # 截断超大表
                    page_number=1,
                    metadata={"source_file": file_path, "sheet": sheet_name, "rows": len(rows)},
                ))

        wb.close()
        logger.info("Excel 提取完成: %d 元素 (%s)", len(elements), Path(file_path).name)
        return elements

    # ------------------------------------------------------------------
    # PowerPoint
    # ------------------------------------------------------------------

    def _extract_pptx(self, file_path: str) -> List[DocumentElement]:
        """提取 PPT：每页幻灯片 → 标题 + 正文 + 备注。"""
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx 未安装。pip install python-pptx")
            return self._fallback_text(file_path)

        prs = Presentation(file_path)
        elements = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                title = texts[0] if texts else ""
                body = " ".join(texts[1:]) if len(texts) > 1 else ""
                if title:
                    elements.append(DocumentElement(
                        category="Title", text=title, page_number=slide_num,
                        metadata={"source_file": file_path},
                    ))
                if body:
                    elements.append(DocumentElement(
                        category="NarrativeText", text=body, page_number=slide_num,
                        metadata={"source_file": file_path},
                    ))

        logger.info("PPT 提取完成: %d 元素 (%s)", len(elements), Path(file_path).name)
        return elements

    # ------------------------------------------------------------------
    # Fallback
    # ------------------------------------------------------------------

    def _fallback_text(self, file_path: str) -> List[DocumentElement]:
        """纯文本 fallback。"""
        try:
            text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
            return [DocumentElement(
                category="NarrativeText", text=text.strip(),
                page_number=1, metadata={"source_file": file_path},
            )]
        except Exception:
            return []
